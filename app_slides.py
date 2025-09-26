import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import requests
from io import BytesIO
import re

# Função para buscar imagem relevante (usando Unsplash API - gratuita, sem chave necessária para buscas básicas)
def buscar_imagem(keyword):
    try:
        url = f"https://source.unsplash.com/800x600/?{keyword}&random"
        response = requests.get(url)
        if response.status_code == 200:
            image = Image.open(BytesIO(response.content))
            return image
        return None
    except:
        return None

# Função para extrair palavras-chave do texto
def extrair_keywords(texto):
    # Simples: pega palavras comuns, remove stop words básicas
    palavras = re.findall(r'\b[a-zA-Z]{4,}\b', texto.lower())
    stop_words = {'o', 'a', 'de', 'do', 'da', 'em', 'para', 'com', 'um', 'uma', 'os', 'as', 'no', 'na', 'é', 'que', 'e', 'ou'}
    keywords = [p for p in palavras if p not in stop_words][:3]  # Top 3 keywords
    return ' '.join(keywords) if keywords else 'education'

# Função principal para gerar PPTX
def gerar_slides(titulo, conteudo):
    prs = Presentation()  # Cria apresentação vazia

    # Slide 1: Título
    slide_layout = prs.slide_layouts[0]  # Layout de título
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = titulo
    subtitle.text = "Trabalho de Escola - Gerado Automaticamente"
    
    # Formatação top: Fonte grande, cor azul
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(44) if shape == title else Pt(24)
                    run.font.color.rgb = RGBColor(0, 51, 102)  # Azul escuro

    # Busca imagem para título
    keyword = extrair_keywords(titulo)
    img = buscar_imagem(keyword)
    if img:
        img_slide = slide.shapes.add_picture(img, Inches(7), Inches(2), width=Inches(3), height=Inches(2))
        # Ajusta posição

    # Divide conteúdo em slides (simples: um slide por parágrafo principal)
    paragrafos = [p.strip() for p in conteudo.split('\n') if p.strip()]
    for i, par in enumerate(paragrafos, 2):
        # Layout de conteúdo com título e texto
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = f"Slide {i}: {par[:50]}..."  # Título curto do parágrafo
        content.text = par
        
        # Formatação: Alinhamento central, fonte sans-serif
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape.text_frame.word_wrap = True
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(18)
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(0, 0, 0)  # Preto

        # Adiciona imagem relevante
        keyword = extrair_keywords(par)
        img = buscar_imagem(keyword)
        if img:
            left = Inches(0.5)
            top = Inches(2.5)
            slide.shapes.add_picture(img, left, top, width=Inches(4), height=Inches(3))

    # Slide final: Conclusão
    slide_layout = prs.slide_layouts[6]  # Layout em branco para custom
    slide = prs.slides.add_slide(slide_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.text = "Obrigado pela atenção!\n\nGerado com IA para facilitar sua vida."
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(32)
            run.font.color.rgb = RGBColor(0, 128, 0)  # Verde

    # Salva o arquivo
    filename = f"slides_{titulo.replace(' ', '_')[:20]}.pptx"
    prs.save(filename)
    return filename

# Interface Streamlit
st.title("🖥️ Gerador de Slides para Trabalho de Escola")
st.write("Digite o título e o conteúdo. O app faz o resto: slides profissionais com imagens! 📚✨")

# Inputs simples
titulo = st.text_input("Título do Trabalho:", placeholder="Ex: O Impacto do Meio Ambiente")
conteudo = st.text_area("Conteúdo (escreva parágrafos separados por Enter):", 
                        placeholder="Ex: Introdução: O meio ambiente é importante...\n\nCorpo: Vamos falar sobre reciclagem...\n\nConclusão: Precisamos agir agora!")

if st.button("🔥 Gerar Slides Top de Linha!"):
    if titulo and conteudo:
        with st.spinner("Criando slides incríveis... Buscando imagens relevantes..."):
            arquivo = gerar_slides(titulo, conteudo)
            st.success(f"Slides gerados! Baixe o arquivo: **{arquivo}**")
            with open(arquivo, "rb") as f:
                st.download_button("📥 Baixar PPTX", f.read(), file_name=arquivo)
    else:
        st.warning("Preencha o título e o conteúdo!")

st.info("Dica: Escreva o conteúdo como um texto normal. O app divide em slides automáticos. Imagens vêm de fontes gratuitas e seguras.")
streamlit run app.py
